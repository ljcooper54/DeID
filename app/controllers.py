# Copyright 2025 H2so4 Consulting LLC

import os
from typing import List
from core.persistence import Persistence  # DB access
from core.entity_detector import EntityDetector  # NLP + heuristics
from core.pseudonym_mapper import PseudonymMapper  # pseudonym generator / reuse
from core.obscure_service import ObscureService  # anonymization pipeline
from core.restore_service import RestoreService  # de-anonymization pipeline
from core import hash_utils  # hashing for file identity

# Document loaders for multi-format support
from docx import Document as DocxDocument            # .docx
import openpyxl                                      # .xlsx
from pptx import Presentation                        # .pptx
from pdfminer.high_level import extract_text         # .pdf


class AppController:
    # AppController: top-level coordinator for login, project state, file association,
    # name lists, multi-format load, obscuring/restoring, and UI-facing info.

    def __init__(self, db_path: str):
        # __init__: create service objects and track current user / project.
        self.db = Persistence(db_path)
        self.detector = EntityDetector()
        self.mapper = PseudonymMapper(self.db)
        self.obscurer = ObscureService(self.detector, self.mapper, self.db)
        self.restorer = RestoreService(self.db)

        self._current_user_id = None
        self._current_project_id = None
        # __init__  # AppController.__init__

    # ------------------------------------------------------------------
    # AUTH
    # ------------------------------------------------------------------

    def create_user(self, username: str, password: str) -> int:
        # create_user: create a brand new user and set them as current.
        uid = self.db.create_user(username, password)
        self._current_user_id = uid
        self._current_project_id = None
        return uid
        # create_user  # AppController.create_user

    def login(self, username: str, password: str) -> bool:
        # login: validate credentials. If OK, set current user
        # and restore their last active project (if valid).
        uid = self.db.validate_login(username, password)
        if uid is None:
            return False

        self._current_user_id = uid

        last_pid = self.db.get_last_project_for_user(uid)
        if last_pid is not None:
            try:
                owner = self.db.get_project_owner(last_pid)
                if owner == uid:
                    self._current_project_id = last_pid
                else:
                    self._current_project_id = None
            except Exception:
                self._current_project_id = None
        else:
            self._current_project_id = None

        return True
        # login  # AppController.login

    def get_current_user_id(self):
        # get_current_user_id: return active user_id or None.
        return self._current_user_id
        # get_current_user_id  # AppController.get_current_user_id

    def get_current_project_id(self):
        # get_current_project_id: return active project_id or None.
        return self._current_project_id
        # get_current_project_id  # AppController.get_current_project_id

    # ------------------------------------------------------------------
    # PROJECTS
    # ------------------------------------------------------------------

    def create_project(self, name: str, notes: str = "") -> int:
        # create_project: create a project for the logged-in user, set it active,
        # and persist it as their "last project".
        if self._current_user_id is None:
            raise RuntimeError("Not logged in.")
        pid = self.db.create_project(self._current_user_id, name, notes)

        self._current_project_id = pid
        self.db.set_last_project_for_user(self._current_user_id, pid)

        return pid
        # create_project  # AppController.create_project

    def select_project(self, project_id: int):
        # select_project: set current active project, validating ownership
        # and persisting last_project_id for this user.
        if self._current_user_id is None:
            raise RuntimeError("Not logged in.")
        owner = self.db.get_project_owner(project_id)
        if owner != self._current_user_id:
            raise RuntimeError("Access denied to that project.")

        self._current_project_id = project_id
        self.db.set_last_project_for_user(self._current_user_id, project_id)
        # select_project  # AppController.select_project

    def list_projects(self):
        # list_projects: return all projects owned by the active user.
        if self._current_user_id is None:
            raise RuntimeError("Not logged in.")
        return self.db.list_projects_for_user(self._current_user_id)
        # list_projects  # AppController.list_projects

    def list_project_files(self):
        # list_project_files: return the files known/associated with the active project.
        if self._current_project_id is None:
            raise RuntimeError("No active project.")
        return self.db.list_project_files(self._current_project_id)
        # list_project_files  # AppController.list_project_files

    def add_files_to_current_project(self, file_paths: List[str]):
        # add_files_to_current_project: add any number of file paths into the active project.
        if self._current_project_id is None:
            raise RuntimeError("No active project.")
        for p in file_paths:
            self.db.upsert_project_file(
                self._current_project_id,
                hash_utils.path_hash(p),
                display_name=p,
                obscured_path=None
            )
        # end for  # file_paths loop
        # add_files_to_current_project  # AppController.add_files_to_current_project

    # ------------------------------------------------------------------
    # FILE LOAD/SAVE HELPERS
    # ------------------------------------------------------------------

    def _load_file_as_text(self, filepath: str) -> str:
        # _load_file_as_text: extract human-readable text from files (.txt, .docx, .xlsx, .pptx, .pdf).
        _, ext = os.path.splitext(filepath)
        ext = ext.lower()

        if ext == ".txt":
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        if ext == ".docx":
            doc = DocxDocument(filepath)
            parts = []
            for para in doc.paragraphs:
                parts.append(para.text)
            # end for
            return "\n".join(parts)

        if ext == ".xlsx":
            wb = openpyxl.load_workbook(filepath, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"=== Sheet: {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [("" if v is None else str(v)) for v in row]
                    parts.append("\t".join(row_vals))
                # end for
            # end for
            return "\n".join(parts)

        if ext == ".pptx":
            pres = Presentation(filepath)
            parts = []
            for slide_idx, slide in enumerate(pres.slides):
                parts.append(f"=== Slide {slide_idx+1} ===")
                for shp in slide.shapes:
                    if hasattr(shp, "text"):
                        text_val = getattr(shp, "text", "")
                        if text_val:
                            parts.append(text_val)
                # end for
            # end for
            return "\n".join(parts)

        if ext == ".pdf":
            return extract_text(filepath)

        # fallback for unknown extension: treat as text
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
        # _load_file_as_text  # AppController._load_file_as_text

    def _write_obscured_output(self,
                               source_path: str,
                               obscured_text: str) -> str:
        # _write_obscured_output: write obscured text to "Obscured_<original-filename>"
        # in the same directory as `source_path`.
        outpath = self.obscurer.build_obscured_filename(source_path)
        with open(outpath, "w", encoding="utf-8") as f:
            f.write(obscured_text)
        return outpath
        # _write_obscured_output  # AppController._write_obscured_output

    def _write_restored_output(self,
                               source_path: str,
                               restored_text: str) -> str:
        # _write_restored_output: write restored text to "Restored_<...>" filename
        # in the same directory as `source_path`.
        outpath = self.restorer.build_restored_filename(source_path)
        with open(outpath, "w", encoding="utf-8") as f:
            f.write(restored_text)
        return outpath
        # _write_restored_output  # AppController._write_restored_output

    # ------------------------------------------------------------------
    # OBFUSCATION / RESTORATION
    # ------------------------------------------------------------------

    def obscure_files(self, file_paths: List[str]) -> List[str]:
        # obscure_files: run anonymization on one or more files under the active project,
        # write Obscured_* outputs, and update project_file rows with the obscured path + timestamp.
        if self._current_project_id is None:
            raise RuntimeError("No active project.")

        outpaths: List[str] = []

        for filepath in file_paths:
            # read file text
            text = self._load_file_as_text(filepath)

            # run obscurer
            out_result = self.obscurer.obscure_text(
                project_id=self._current_project_id,
                file_path=filepath,
                text=text
            )

            # write output file using Obscured_ naming
            out_path = self._write_obscured_output(filepath, out_result.obscured_text)
            outpaths.append(out_path)

            # remember this file in project_file, and note where the obscured version lives
            self.db.update_project_file_after_obscure(
                self._current_project_id,
                source_display_name=filepath,
                source_hash=hash_utils.path_hash(filepath),
                new_obscured_path=out_path
            )
        # end for  # file_paths loop

        return outpaths
        # obscure_files  # AppController.obscure_files

    def restore_files(self, file_paths: List[str]) -> List[str]:
        # restore_files: run de-anonymization on one or more obscured text files
        # under the active project, and save as Restored_*.
        if self._current_project_id is None:
            raise RuntimeError("No active project.")

        outpaths: List[str] = []

        for filepath in file_paths:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                obscured_text = f.read()

            res = self.restorer.restore_text(self._current_project_id, obscured_text)

            out_path = self._write_restored_output(filepath, res.restored_text)
            outpaths.append(out_path)
        # end for  # file_paths loop

        return outpaths
        # restore_files  # AppController.restore_files

    # ------------------------------------------------------------------
    # GLOBAL / PROJECT SENSITIVE NAMES
    # ------------------------------------------------------------------

    def import_user_names_list(self, path_to_file: str):
        # import_user_names_list: import newline-delimited strings into the user's global list.
        if self._current_user_id is None:
            raise RuntimeError("Not logged in.")
        with open(path_to_file, "r", encoding="utf-8") as f:
            for line in f:
                candidate = line.strip()
                if not candidate:
                    continue
                self.db.add_user_known_name(self._current_user_id, candidate)
        # import_user_names_list  # AppController.import_user_names_list

    def import_project_names_list(self, path_to_file: str):
        # import_project_names_list: import newline-delimited strings into the active project's list.
        if self._current_project_id is None:
            raise RuntimeError("No active project.")
        with open(path_to_file, "r", encoding="utf-8") as f:
            for line in f:
                candidate = line.strip()
                if not candidate:
                    continue
                self.db.add_project_known_name(self._current_project_id, candidate)
        # import_project_names_list  # AppController.import_project_names_list

    def import_user_names_list_from_values(self, values: List[str]):
        # import_user_names_list_from_values: add provided strings directly to the user's global list.
        if self._current_user_id is None:
            raise RuntimeError("Not logged in.")
        for v in values:
            if not v.strip():
                continue
            self.db.add_user_known_name(self._current_user_id, v.strip())
        # import_user_names_list_from_values  # AppController.import_user_names_list_from_values

    def import_project_names_list_from_values(self, values: List[str]):
        # import_project_names_list_from_values: add provided strings to active project's forced list.
        if self._current_project_id is None:
            raise RuntimeError("No active project.")
        for v in values:
            if not v.strip():
                continue
            self.db.add_project_known_name(self._current_project_id, v.strip())
        # import_project_names_list_from_values  # AppController.import_project_names_list_from_values

    def list_user_names(self):
        # list_user_names: return the current user's global forced-redaction list.
        if self._current_user_id is None:
            raise RuntimeError("Not logged in.")
        return self.db.list_user_known_names(self._current_user_id)
        # list_user_names  # AppController.list_user_names

    def list_project_names(self):
        # list_project_names: return the active project's forced-redaction list.
        if self._current_project_id is None:
            raise RuntimeError("No active project.")
        return self.db.list_project_known_names(self._current_project_id)
        # list_project_names  # AppController.list_project_names

    def delete_user_name(self, name_text: str):
        # delete_user_name: remove one string from the current user's global forced list.
        if self._current_user_id is None:
            raise RuntimeError("Not logged in.")
        self.db.delete_user_known_name(self._current_user_id, name_text)
        # delete_user_name  # AppController.delete_user_name

    def delete_project_name(self, name_text: str):
        # delete_project_name: remove one string from the active project's forced list.
        if self._current_project_id is None:
            raise RuntimeError("No active project.")
        self.db.delete_project_known_name(self._current_project_id, name_text)
        # delete_project_name  # AppController.delete_project_name

# AppController
